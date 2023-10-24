# https://python-pptx.readthedocs.io/en/latest/

from pptx import Presentation

def shift_text(ppt_name, old, new):
    """
    https://stackoverflow.com/questions/50962836/replacing-particular-text-in-all-sides-of-a-ppt-using-python-pptx

    > Extras
    * https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
      * python -m pip install python-pptx-text-replacer
      * python -m pip install petpptx
    * https://stackoverflow.com/questions/55497789/find-a-word-in-multiple-powerpoint-files-python
    """
    prs = Persentation(ppt_name)
    for slide in prs.slides: 
        for shape in slide.shapes: 
            if not shape.has_text_frame: 
                    continue 
            for paragraph in shape.text_frame.paragraphs:  
                for run in paragraph.runs:
                    run.text=newText(run.text)  
    prs.save(ppt_name) 

#tkinter
#https://realpython.com/python-gui-tkinter/

import tkinter as tk

window = tk.Tk()
label = tk.Label(text="Python rocks!")
label.pack()

window.mainloop()
